from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_FILE = Path(__file__).resolve().parent.parent / "滑雪装备多站点库存与闲鱼分发系统.pptx"

BG = RGBColor(244, 241, 235)
PANEL = RGBColor(255, 252, 246)
PRIMARY = RGBColor(26, 59, 93)
ACCENT = RGBColor(204, 86, 38)
TEXT = RGBColor(35, 35, 35)
MUTED = RGBColor(106, 106, 106)
LINE = RGBColor(221, 214, 201)


SLIDES = [
    {
        "type": "cover",
        "title": "滑雪装备多站点库存与闲鱼分发系统",
        "subtitle": "从韩国站点抓取、库存同步、定价规则到闲鱼多账号发布的一体化运营后台",
        "meta": "Boardline Spider / Xianyu Open System   2026",
    },
    {
        "title": "项目定位",
        "bullets": [
            "这不是单点脚本，而是一套完整的自动化运营系统。",
            "对接韩国滑雪装备网站，持续抓取商品、图片、价格、库存。",
            "把原始商品数据整理成可运营、可定价、可发布的数据资产。",
            "支撑闲鱼多账号、多批次、多规格商品的自动化分发。",
        ],
        "footer": "核心目标：把分散在网站、表格、人工操作里的流程统一到一个系统里。",
    },
    {
        "title": "解决的核心问题",
        "bullets": [
            "商品来源分散，人工查看与更新慢，容易漏品。",
            "库存变化频繁，热门型号与尺码靠人工跟踪成本高。",
            "发布链路复杂，标题、简介、图片、规格、类目都要人工处理。",
            "多账号管理困难，创建、上架、失败、回调难以追踪。",
        ],
        "footer": "结果：上新慢、库存不同步、重复劳动多、发布结果不透明。",
    },
    {
        "title": "系统整体架构",
        "bullets": [
            "韩国站点分类抓取",
            "商品列表抓取",
            "商品库存与价格解析",
            "本地数据库沉淀",
            "规则化定价与折扣",
            "AI标题 / 简介 / 图片生成",
            "闲鱼批次创建与发布",
            "回调、状态同步、失败重试",
        ],
        "footer": "形成“抓取端 + 数据端 + 运营端 + 发布端”的业务闭环。",
    },
    {
        "title": "多站点抓取能力",
        "bullets": [
            "已成熟运行：Boardline。",
            "已完成接入：One8。",
            "通过 source + branduid 区分不同来源网站。",
            "同站点重跑会更新已有商品，新商品会自动入库。",
            "为后续扩展 4-6 个韩国站点预留了统一结构。",
        ],
        "footer": "不是一次性抓取，而是可持续更新的多来源商品底座。",
    },
    {
        "title": "库存能力",
        "bullets": [
            "支持单规格库存、多规格库存、套装商品库存。",
            "支持附加选项清洗后归并主规格。",
            "支持 0 库存过滤与“售罄”标记。",
            "支持库存异常补抓与失败修复。",
            "库存结构可直接用于后续多规格商品发布。",
        ],
        "footer": "库存不是简单文本，而是可继续分发的结构化数据。",
    },
    {
        "title": "价格与利润规则",
        "bullets": [
            "Boardline 已有定价规则表与折扣规则表。",
            "One8 已新增独立 pricing_rules 与 discount_rules。",
            "支持按品类设置汇率、运费、利润率。",
            "支持按关键词设置折扣类型、折扣值与优先级。",
            "运营人员可以直接改表，无需改代码。",
        ],
        "footer": "不同来源网站可以独立定价，价格策略更灵活。",
    },
    {
        "title": "闲鱼分发能力",
        "bullets": [
            "支持多闲鱼账号管理。",
            "支持商品池筛选、勾选、生成批次。",
            "支持批量创建、批量上架、批量删除。",
            "支持失败重试、待发布重处理、查询后重发。",
            "支持回调状态同步与异步任务进度显示。",
        ],
        "footer": "把单次上传动作升级成可追踪、可重试、可审计的任务系统。",
    },
    {
        "title": "AI能力",
        "bullets": [
            "支持 AI 标题生成、AI 简介生成、AI 主图生成。",
            "新账号支持独立 AI 内容，不与旧账号共用。",
            "图片生成可按账号定制风格。",
            "创建时可选择是否带水印。",
        ],
        "footer": "降低内容同质化，提升发布效率和账号差异化运营能力。",
    },
    {
        "title": "多账号差异化运营",
        "bullets": [
            "新账号默认独立 AI 文案和图片。",
            "新账号默认配置 10 组发货地区。",
            "每 10 个商品自动轮换地区。",
            "支持账号级发布参数与类目映射。",
        ],
        "footer": "不是机械的一库多发，而是同源商品的多账号差异化运营。",
    },
    {
        "title": "后台能力",
        "bullets": [
            "账号概览、商品池浏览、分类筛选。",
            "批次管理、任务明细查看、回调状态查看。",
            "创建 / 上架 / 删除进度可视化。",
            "失败项可继续补处理。",
        ],
        "footer": "运营动作可视化，结果可回查，异常可定位。",
    },
    {
        "title": "数据资产沉淀",
        "bullets": [
            "商品主表。",
            "商品价格与库存表。",
            "变化日志。",
            "闲鱼账号表、批次表、任务表。",
            "全量库存导出表与每日变化表。",
        ],
        "footer": "业务不再停留在人工层，而是在持续形成自己的商品数据库和运营数据库。",
    },
    {
        "title": "系统卖点",
        "bullets": [
            "多站点商品采集能力。",
            "结构化库存解析能力。",
            "表格化价格与折扣规则。",
            "多账号闲鱼分发能力。",
            "AI 内容差异化能力。",
            "批次与任务化运营能力。",
            "自动补抓与失败修复能力。",
        ],
        "footer": "这不是一个抓取脚本，而是一套面向实际运营的自动化商品分发系统。",
    },
    {
        "title": "业务价值",
        "bullets": [
            "更快上新，减少人工查品与整理时间。",
            "更准库存，降低超卖、误卖、错发风险。",
            "更强复用，同一商品资产可复用到多个账号。",
            "更低人力，规则化处理替代重复劳动。",
            "更易扩展，后续继续接站点、接平台、接账号都更快。",
        ],
        "footer": "适合韩国滑雪装备分销、多平台铺货、多账号精细化运营。",
    },
    {
        "title": "后续升级方向",
        "bullets": [
            "接入更多韩国站点。",
            "继续提升抓取速度，减少浏览器依赖。",
            "增加登录与权限系统。",
            "支持淘宝 / 其他平台分发。",
            "增强价格与库存预警。",
            "提升 AI 图与文案策略能力。",
        ],
        "footer": "当前系统已经可用，后续重点是从“可用”走向“规模化”。",
    },
    {
        "title": "收尾页",
        "bullets": [
            "数据抓取只是入口。",
            "规则系统是中台。",
            "多账号发布是业务出口。",
            "系统的价值，在于把商品数据真正变成运营能力。",
        ],
        "footer": "谢谢观看",
    },
]


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_band(slide):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.55))
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY
    band.line.fill.background()


def add_title(slide, title):
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.8), Inches(11.6), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    r.font.name = "Arial"


def add_footer(slide, text):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(6.55), Inches(11.8), Inches(0.5))
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = LINE
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(11)
    r.font.color.rgb = MUTED
    r.font.name = "Arial"


def add_bullets(slide, bullets):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.55), Inches(11.8), Inches(4.7))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = LINE
    tf = panel.text_frame
    tf.clear()
    tf.word_wrap = True

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(12)
        p.bullet = True
        r = p.add_run()
        r.text = bullet
        r.font.size = Pt(21 if len(bullets) <= 5 else 18)
        r.font.color.rgb = TEXT
        r.font.name = "Arial"


def add_cover(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(4.6), Inches(7.5))
    left.fill.solid()
    left.fill.fore_color.rgb = PRIMARY
    left.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(4.2), Inches(0.7), Inches(0.22), Inches(5.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(4.8), Inches(1.1), Inches(7.3), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = slide_data["title"]
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    r.font.name = "Arial"

    sub_box = slide.shapes.add_textbox(Inches(4.85), Inches(3.55), Inches(6.8), Inches(1.4))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = slide_data["subtitle"]
    r.font.size = Pt(16)
    r.font.color.rgb = MUTED
    r.font.name = "Arial"

    meta_box = slide.shapes.add_textbox(Inches(4.85), Inches(6.4), Inches(6.5), Inches(0.4))
    p = meta_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = slide_data["meta"]
    r.font.size = Pt(12)
    r.font.color.rgb = ACCENT
    r.font.bold = True
    r.font.name = "Arial"


def add_content_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    add_title(slide, slide_data["title"])
    add_bullets(slide, slide_data["bullets"])
    add_footer(slide, slide_data["footer"])


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in SLIDES:
        if slide_data.get("type") == "cover":
            add_cover(prs, slide_data)
        else:
            add_content_slide(prs, slide_data)

    prs.save(str(OUT_FILE))
    print(OUT_FILE)


if __name__ == "__main__":
    main()
